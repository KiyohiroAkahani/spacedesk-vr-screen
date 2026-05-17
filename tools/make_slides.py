#!/usr/bin/env python3
"""Generate dist/SETUP_GUIDE.pptx from the SETUP_GUIDE content.

Trilingual (EN / 中文 / 日本語), beginner-friendly, with fade slide
transitions. Re-run this whenever the setup steps change.

Usage:
    py -3 .\\tools\\make_slides.py

Requires the `python-pptx` package; it is auto-installed on first run if
missing (pip).

Exit codes:
    0  pptx created (path printed)
    1  generation failed
    2  setup error (python-pptx missing and auto-install failed)
"""

from __future__ import annotations

import datetime as _dt
import subprocess
import sys
from pathlib import Path

REPO_ROOT = Path(__file__).resolve().parent.parent
DIST = REPO_ROOT / "dist"
TMP = REPO_ROOT / ".tmp"


def _ensure_pptx() -> None:
    try:
        import pptx  # noqa: F401
        return
    except ImportError:
        pass
    print("[make_slides] 'python-pptx' not found — installing via pip...")
    r = subprocess.run([sys.executable, "-m", "pip", "install",
                        "--disable-pip-version-check", "python-pptx"])
    if r.returncode != 0:
        print("ERROR: failed to install python-pptx. Run:\n"
              "  py -3 -m pip install python-pptx", file=sys.stderr)
        raise SystemExit(2)
    try:
        import pptx  # noqa: F401
    except ImportError:
        print("ERROR: python-pptx still unavailable.", file=sys.stderr)
        raise SystemExit(2)


def build() -> Path:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.oxml.ns import qn
    from pptx.oxml import parse_xml

    GREEN = RGBColor(0x3F, 0xB2, 0x3F)
    DARK = RGBColor(0x1E, 0x1E, 0x1E)
    GREY = RGBColor(0x55, 0x55, 0x55)
    RED = RGBColor(0xCC, 0x00, 0x00)
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    LIGHT = RGBColor(0xF4, 0xF7, 0xF4)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    BLANK = prs.slide_layouts[6]
    SW, SH = prs.slide_width, prs.slide_height

    def add_fade(slide):
        xml = ('<p:transition xmlns:p="http://schemas.openxmlformats.org/'
               'presentationml/2006/main" spd="med" advClick="1">'
               '<p:fade/></p:transition>')
        el = parse_xml(xml)
        sld = slide._element
        cz = sld.find(qn('p:clrMapOvr'))
        if cz is not None:
            cz.addnext(el)
        else:
            sld.append(el)

    def box(slide, x, y, w, h, fill=None, line=None, round_=False):
        sp = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE if round_ else MSO_SHAPE.RECTANGLE,
            x, y, w, h)
        sp.shadow.inherit = False
        if fill is None:
            sp.fill.background()
        else:
            sp.fill.solid()
            sp.fill.fore_color.rgb = fill
        if line is None:
            sp.line.fill.background()
        else:
            sp.line.color.rgb = line
            sp.line.width = Pt(2.5)
        return sp

    def text(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.TOP):
        tb = slide.shapes.add_textbox(x, y, w, h)
        tf = tb.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = anchor
        for i, (sstr, sz, color, bold) in enumerate(runs):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = align
            p.space_after = Pt(6)
            r = p.add_run()
            r.text = sstr
            r.font.size = Pt(sz)
            r.font.bold = bold
            r.font.color.rgb = color
            r.font.name = "Segoe UI"
        return tb

    def base(step, title, band=GREEN):
        s = prs.slides.add_slide(BLANK)
        box(s, 0, 0, SW, SH, fill=WHITE)
        box(s, 0, 0, SW, Inches(1.15), fill=band)
        if step is not None:
            b = box(s, Inches(0.45), Inches(0.16), Inches(0.83),
                    Inches(0.83), fill=WHITE, round_=True)
            tf = b.text_frame
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            r = p.add_run()
            r.text = str(step)
            r.font.size = Pt(34)
            r.font.bold = True
            r.font.color.rgb = band
        text(s, Inches(1.45 if step is not None else 0.5), Inches(0.18),
             Inches(11.4), Inches(0.8), [(title, 30, WHITE, True)],
             anchor=MSO_ANCHOR.MIDDLE)
        add_fade(s)
        return s

    def tri(slide, y, en, zh, ja, h=Inches(1.5), big=False):
        box(slide, Inches(0.6), y, Inches(12.13), h, fill=LIGHT, round_=True)
        fs = 22 if big else 18
        text(slide, Inches(0.95), y + Inches(0.12), Inches(11.5),
             h - Inches(0.24),
             [("EN  " + en, fs, DARK, big),
              ("中文  " + zh, fs, GREY, False),
              ("日本語  " + ja, fs, DARK, big)])

    # 1 Title
    s = prs.slides.add_slide(BLANK)
    box(s, 0, 0, SW, SH, fill=DARK)
    box(s, 0, Inches(2.55), SW, Inches(0.08), fill=GREEN)
    text(s, Inches(0.8), Inches(1.5), Inches(11.7), Inches(1.1),
         [("spacedesk-vr-screen", 48, WHITE, True)], align=PP_ALIGN.CENTER)
    text(s, Inches(0.8), Inches(2.75), Inches(11.7), Inches(2.4),
         [("Setup Guide — turn your phone into a VR screen of your PC",
           24, GREEN, True),
          ("安装指南 / セットアップガイド  (English / 中文 / 日本語)",
           20, WHITE, False),
          ("Follow the steps in order — Step 3 (display) is the key one.",
           18, RGBColor(0xCC, 0xCC, 0xCC), False)], align=PP_ALIGN.CENTER)
    add_fade(s)

    # 2 What you need
    s = base("0", "What you need / 需要准备 / 必要なもの")
    tri(s, Inches(1.5),
        "Windows 10/11 PC + a smartphone (iPhone/Android).",
        "Windows 10/11 电脑 + 一部智能手机（iPhone/Android）。",
        "Windows 10/11 のPC ＋ スマホ（iPhone/Android）。")
    tri(s, Inches(3.15),
        "A phone VR goggle that fits your phone (see Step 4).",
        "一个适配你手机的手机 VR 头显（见第 4 步）。",
        "自分のスマホに合うVRゴーグル（手順4）。")
    tri(s, Inches(4.8),
        "PC and phone able to join the SAME Wi-Fi.",
        "电脑和手机能连到同一个 Wi-Fi。",
        "PCとスマホを同じWi-Fiにつなげること。")

    # 3 Step 1 install
    s = base("1", "Install spacedesk (2 parts) / 安装 spacedesk / spacedeskを入れる")
    tri(s, Inches(1.45),
        "PC: from spacedesk.net download 'spacedesk DRIVER for Windows "
        "(Server)' and install it.",
        "电脑：从 spacedesk.net 下载 'spacedesk DRIVER for Windows "
        "(Server)' 并安装。",
        "PC：spacedesk.net から『spacedesk DRIVER for Windows (Server)』"
        "を導入。", h=Inches(1.7))
    box(s, Inches(0.6), Inches(3.35), Inches(12.13), Inches(1.05),
        fill=RGBColor(0xFF, 0xEC, 0xEC), line=RED, round_=True)
    text(s, Inches(0.95), Inches(3.5), Inches(11.5), Inches(0.8),
         [("⚠  NOT the Microsoft Store 'spacedesk' app  /  不是商店版  /  "
           "Microsoft Store版は使わない", 19, RED, True)],
         anchor=MSO_ANCHOR.MIDDLE)
    tri(s, Inches(4.65),
        "Phone: install 'spacedesk - USB display for PC' (App Store / "
        "Google Play).",
        "手机：安装 'spacedesk - USB display for PC'（App Store / "
        "Google Play）。",
        "スマホ：『spacedesk - USB display for PC』を導入。", h=Inches(1.7))

    # 4 Step 2 wifi + app
    s = base("2", "Same Wi-Fi & get the app / 同一Wi-Fi与应用 / 同じWi-Fiとアプリ")
    tri(s, Inches(1.45),
        "Put the PC and phone on the SAME Wi-Fi.",
        "让电脑和手机连到同一个 Wi-Fi。",
        "PCとスマホを同じWi-Fiに接続。")
    tri(s, Inches(3.05),
        "GitHub repo -> 'Releases' -> download VrDesktopBridge.exe.",
        "GitHub 仓库 -> 'Releases' -> 下载 VrDesktopBridge.exe。",
        "GitHubの『Releases』-> VrDesktopBridge.exe をダウンロード。",
        h=Inches(1.55))
    box(s, Inches(0.6), Inches(4.75), Inches(12.13), Inches(1.5),
        fill=RGBColor(0xFF, 0xF6, 0xE5), line=RGBColor(0xE0, 0x90, 0x00),
        round_=True)
    text(s, Inches(0.95), Inches(4.9), Inches(11.5), Inches(1.2),
         [("ℹ  First run: 'Windows protected your PC'",
           18, RGBColor(0xB0, 0x70, 0x00), True),
          ("-> 'More info' -> 'Run anyway'  /  '更多信息'->'仍要运行'  /  "
           "『詳細情報』->『実行』", 17, DARK, False)])

    # 5 Step 3a connect
    s = base("3", "Connect / 连接 / 接続")
    tri(s, Inches(1.6),
        "On the PC: start spacedesk (Console), leave it running.",
        "电脑：启动 spacedesk（Console），保持运行。",
        "PC：spacedesk（Console）を起動したままにする。", h=Inches(1.7))
    tri(s, Inches(3.5),
        "On the phone: open spacedesk -> tap your PC name to connect.",
        "手机：打开 spacedesk -> 点你的电脑名连接。",
        "スマホ：spacedesk を開き PC名をタップして接続。", h=Inches(1.9))

    # 6 Step 3b EXTEND (key)
    s = base("3", "Set the display to 'Extend' — KEY STEP", band=RED)
    box(s, Inches(0.6), Inches(1.45), Inches(12.13), Inches(2.05),
        fill=RGBColor(0xFF, 0xEC, 0xEC), line=RED, round_=True)
    text(s, Inches(0.95), Inches(1.62), Inches(11.5), Inches(1.8),
         [("Right-click desktop -> 'Display settings' -> 'Multiple "
           "displays'", 20, DARK, True),
          ("-> choose  OK 'Extend these displays'", 22, GREEN, True),
          ("X  NOT 'Duplicate these displays'  (Duplicate = cannot "
           "operate, app fails)", 19, RED, True)])
    tri(s, Inches(3.75),
        "Exact: 'Extend these displays' (NOT 'Duplicate').",
        "原文：'扩展这些显示器'（不是 '复制这些显示器'）。",
        "画面の文言：『表示画面を拡張する』（『複製』は不可）。",
        h=Inches(1.45), big=True)
    text(s, Inches(0.6), Inches(5.45), Inches(12.13), Inches(0.7),
         [("Quick / 快捷 / 簡単:  Windows key + P  ->  'Extend' / '扩展' "
           "/ 『拡張』", 18, GREY, True)], align=PP_ALIGN.CENTER)

    # 7 Step 4 goggle + launch
    s = base("4", "VR goggle & launch / 头显与启动 / ゴーグルと起動")
    tri(s, Inches(1.45),
        "Recommended VR goggle (May 2026): SHREVNI on Amazon JP "
        "-> amazon.co.jp/dp/B0FP5M4L6B/",
        "推荐手机 VR 头显（2026/5）：Amazon 日本 SHREVNI "
        "-> amazon.co.jp/dp/B0FP5M4L6B/",
        "おすすめVRゴーグル（2026/5）：Amazon の SHREVNI "
        "-> amazon.co.jp/dp/B0FP5M4L6B/", h=Inches(1.95))
    tri(s, Inches(3.55),
        "Double-click VrDesktopBridge.exe -> side-by-side VR. Put phone "
        "in the goggle.",
        "双击 VrDesktopBridge.exe -> 左右并排 VR。把手机放入头显。",
        "VrDesktopBridge.exe をダブルクリック -> 左右2分割VR。"
        "ゴーグルに装着。", h=Inches(1.95))

    # 8 Step 5 controls
    s = base("5", "Everyday controls / 常用操作 / ふだんの操作")
    tri(s, Inches(1.45),
        "Quit: tap Esc TWICE quickly  (or Ctrl+Alt+Shift+Q).",
        "退出：快速按两次 Esc（或 Ctrl+Alt+Shift+Q）。",
        "終了：Esc を素早く2回（または Ctrl+Alt+Shift+Q）。", big=True)
    tri(s, Inches(3.1),
        "Zoom in / out: Ctrl+Alt+Shift+ Up / Down.",
        "缩放：Ctrl+Alt+Shift+ 上 / 下。",
        "拡大/縮小：Ctrl+Alt+Shift+ ↑ / ↓。", big=True)
    tri(s, Inches(4.75),
        "Lost a window? It is auto-returned (toggle Ctrl+Alt+Shift+W).",
        "窗口不见了？会自动拉回（Ctrl+Alt+Shift+W 切换）。",
        "ウィンドウが見当たらない？自動で戻ります（Ctrl+Alt+Shift+W）。")

    # 9 Step 6 recommended: 4K video  (NEW)
    s = base("6", "Recommended: a 4K video / 推荐：4K 视频 / おすすめ：4K動画",
             band=GREEN)
    tri(s, Inches(1.6),
        "For a great first impression, play a high-resolution (4K) "
        "video FULL-SCREEN and watch it on the VR big screen.",
        "想获得惊艳的第一印象，全屏播放高分辨率（4K）视频，在 VR 大屏幕"
        "上观看。",
        "最初の感動を味わうなら、高解像度（4K）の動画を全画面で再生し、"
        "VRの大画面で観るのがおすすめ。", h=Inches(2.1))
    box(s, Inches(0.6), Inches(4.0), Inches(12.13), Inches(1.2),
        fill=RGBColor(0xEC, 0xF7, 0xEC), line=GREEN, round_=True)
    text(s, Inches(0.95), Inches(4.2), Inches(11.5), Inches(0.9),
         [("Example / 示例 / 例:  youtube.com/watch?v=Pt_2nw6vv1k&t=2589s",
           18, DARK, True)], anchor=MSO_ANCHOR.MIDDLE)

    # 10 Troubleshoot
    s = base("?", "If it doesn't work / 出问题时 / うまくいかないとき")
    tri(s, Inches(1.7),
        "Most issues = display mode. Re-check Step 3: 'Extend these "
        "displays' (NOT Duplicate).",
        "多数问题 = 显示模式。重做第 3 步：'扩展这些显示器'（不是复制）。",
        "多くは表示モードが原因。手順3を再確認：『表示画面を拡張する』"
        "（複製は不可）。", h=Inches(1.9))
    tri(s, Inches(3.75),
        "Phone blank? Reconnect spacedesk on the phone, then press "
        "Ctrl+Alt+Shift+D in the app.",
        "手机无画面？手机端重连 spacedesk，再在程序内按 Ctrl+Alt+Shift+D。",
        "スマホが映らない？spacedesk を繋ぎ直し、アプリ内で "
        "Ctrl+Alt+Shift+D。", h=Inches(1.9))

    # 11 Links / thanks
    s = prs.slides.add_slide(BLANK)
    box(s, 0, 0, SW, SH, fill=DARK)
    box(s, 0, Inches(3.05), SW, Inches(0.06), fill=GREEN)
    text(s, Inches(0.8), Inches(1.2), Inches(11.7), Inches(1.0),
         [("Enjoy your desktop in VR!", 40, WHITE, True)],
         align=PP_ALIGN.CENTER)
    text(s, Inches(0.8), Inches(3.3), Inches(11.7), Inches(2.6),
         [("GitHub: github.com/KiyohiroAkahani/spacedesk-vr-screen",
           22, GREEN, True),
          ("Details & all languages: see README / 详见 README / "
           "詳細は README", 18, WHITE, False),
          ("Thanks: Claude Opus, the giants of programming, and the "
           "first tester — my wife.", 16, RGBColor(0xCC, 0xCC, 0xCC),
           False)], align=PP_ALIGN.CENTER)
    add_fade(s)

    DIST.mkdir(exist_ok=True)
    out = DIST / "SETUP_GUIDE.pptx"
    prs.save(str(out))
    return out


def main() -> None:
    _ensure_pptx()
    TMP.mkdir(exist_ok=True)
    stamp = _dt.datetime.now().strftime("%Y%m%d_%H%M%S")
    try:
        out = build()
        from pptx import Presentation
        n = len(Presentation(str(out)).slides)  # reopen = validity check
    except SystemExit:
        raise
    except Exception as ex:  # noqa: BLE001
        log = TMP / f"slides_{stamp}.log"
        log.write_text(repr(ex), encoding="utf-8")
        print(f"ERROR: slide generation failed: {ex}\n(see {log})",
              file=sys.stderr)
        raise SystemExit(1)

    rel = out.relative_to(REPO_ROOT)
    print(f"SLIDES CREATED: {rel}  ({n} slides, "
          f"{out.stat().st_size} bytes)")
    print("Share: attach to a GitHub Release, or open/edit in PowerPoint.")
    raise SystemExit(0)


if __name__ == "__main__":
    main()
